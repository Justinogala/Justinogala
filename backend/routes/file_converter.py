"""
File Converter API — handles format conversions between PDF, Word, Image, Excel, PPTX, EPUB, MOBI.
"""
import io
import os
import subprocess
import tempfile
import zipfile
import logging
from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse
from typing import List
import fitz  # PyMuPDF
from PIL import Image

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/converter", tags=["converter"])

ALLOWED_CONVERSIONS = {
    "pdf-to-jpg", "pdf-to-png", "pdf-to-word",
    "jpg-to-pdf", "png-to-pdf", "image-to-pdf",
    "word-to-pdf", "excel-to-pdf", "pptx-to-pdf",
    "png-to-jpg", "jpg-to-png",
    "epub-to-mobi", "mobi-to-epub", "epub-to-pdf",
}

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB


@router.get("/supported")
async def get_supported_conversions():
    return {
        "conversions": [
            {"id": "pdf-to-jpg", "from": "PDF", "to": "JPG", "category": "Convert from PDF"},
            {"id": "pdf-to-png", "from": "PDF", "to": "PNG", "category": "Convert from PDF"},
            {"id": "pdf-to-word", "from": "PDF", "to": "DOCX", "category": "Convert from PDF"},
            {"id": "word-to-pdf", "from": "DOCX", "to": "PDF", "category": "Convert to PDF"},
            {"id": "excel-to-pdf", "from": "XLSX", "to": "PDF", "category": "Convert to PDF"},
            {"id": "pptx-to-pdf", "from": "PPTX", "to": "PDF", "category": "Convert to PDF"},
            {"id": "jpg-to-pdf", "from": "JPG", "to": "PDF", "category": "Convert to PDF"},
            {"id": "png-to-pdf", "from": "PNG", "to": "PDF", "category": "Convert to PDF"},
            {"id": "image-to-pdf", "from": "Image", "to": "PDF", "category": "Convert to PDF"},
            {"id": "png-to-jpg", "from": "PNG", "to": "JPG", "category": "Image Converter"},
            {"id": "jpg-to-png", "from": "JPG", "to": "PNG", "category": "Image Converter"},
            {"id": "epub-to-mobi", "from": "EPUB", "to": "MOBI", "category": "eBook"},
            {"id": "mobi-to-epub", "from": "MOBI", "to": "EPUB", "category": "eBook"},
            {"id": "epub-to-pdf", "from": "EPUB", "to": "PDF", "category": "eBook"},
        ]
    }


@router.post("/convert")
async def convert_file(
    file: UploadFile = File(...),
    conversion_type: str = Form(...),
):
    if conversion_type not in ALLOWED_CONVERSIONS:
        raise HTTPException(status_code=400, detail=f"Unsupported conversion: {conversion_type}")

    file_data = await file.read()
    if len(file_data) > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail="File too large (max 50MB)")

    original_name = os.path.splitext(file.filename or "file")[0]

    try:
        if conversion_type == "pdf-to-jpg":
            return _pdf_to_images(file_data, original_name, "jpeg")
        elif conversion_type == "pdf-to-png":
            return _pdf_to_images(file_data, original_name, "png")
        elif conversion_type == "pdf-to-word":
            return _pdf_to_word(file_data, original_name)
        elif conversion_type in ("jpg-to-pdf", "png-to-pdf", "image-to-pdf"):
            return _image_to_pdf(file_data, original_name)
        elif conversion_type == "word-to-pdf":
            return _word_to_pdf(file_data, original_name)
        elif conversion_type == "excel-to-pdf":
            return _excel_to_pdf(file_data, original_name)
        elif conversion_type == "pptx-to-pdf":
            return _pptx_to_pdf(file_data, original_name)
        elif conversion_type == "png-to-jpg":
            return _convert_image_format(file_data, original_name, "jpeg")
        elif conversion_type == "jpg-to-png":
            return _convert_image_format(file_data, original_name, "png")
        elif conversion_type == "epub-to-mobi":
            return _calibre_convert(file_data, original_name, ".epub", ".mobi", "application/x-mobipocket-ebook")
        elif conversion_type == "mobi-to-epub":
            return _calibre_convert(file_data, original_name, ".mobi", ".epub", "application/epub+zip")
        elif conversion_type == "epub-to-pdf":
            return _epub_to_pdf(file_data, original_name)
        else:
            raise HTTPException(status_code=400, detail="Conversion not implemented")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Conversion error ({conversion_type}): {e}")
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")


# ── Batch image merge helpers ─────────────────────────────────
MERGE_TO_PDF_TYPES = {"jpg-to-pdf", "png-to-pdf", "image-to-pdf"}


@router.post("/batch-convert")
async def batch_convert(
    files: List[UploadFile] = File(...),
    conversion_type: str = Form(...),
):
    """Convert multiple files at once. Images→PDF merges into one PDF. Others zip individually."""
    if conversion_type not in ALLOWED_CONVERSIONS:
        raise HTTPException(status_code=400, detail=f"Unsupported conversion: {conversion_type}")
    if len(files) < 1:
        raise HTTPException(status_code=400, detail="No files provided")
    if len(files) > 50:
        raise HTTPException(status_code=400, detail="Max 50 files per batch")

    total_size = 0
    file_buffers = []
    for f in files:
        data = await f.read()
        total_size += len(data)
        if total_size > 100 * 1024 * 1024:  # 100MB total
            raise HTTPException(status_code=400, detail="Total batch size exceeds 100MB")
        file_buffers.append((f.filename or "file", data))

    try:
        # Images → combined PDF
        if conversion_type in MERGE_TO_PDF_TYPES:
            return _batch_images_to_pdf(file_buffers)

        # All other conversions → convert each, return ZIP
        return await _batch_individual(file_buffers, conversion_type)
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Batch conversion error ({conversion_type}): {e}")
        raise HTTPException(status_code=500, detail=f"Batch conversion failed: {str(e)}")


def _batch_images_to_pdf(file_buffers):
    """Merge multiple images into a single PDF."""
    pdf = fitz.open()
    for fname, img_data in file_buffers:
        try:
            img = fitz.open(stream=img_data, filetype="png")  # auto-detects
            pdfbytes = img.convert_to_pdf()
            img_pdf = fitz.open("pdf", pdfbytes)
            pdf.insert_pdf(img_pdf)
            img.close()
            img_pdf.close()
        except Exception as e:
            logger.warning(f"Skipping {fname}: {e}")
    if len(pdf) == 0:
        raise HTTPException(status_code=400, detail="No valid images found")
    pdf_bytes = pdf.tobytes()
    pdf.close()
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": 'attachment; filename="combined.pdf"'},
    )


async def _batch_individual(file_buffers, conversion_type):
    """Convert each file individually and return a ZIP."""
    buf = io.BytesIO()
    ext_map = {
        "pdf-to-jpg": "jpg", "pdf-to-png": "png", "pdf-to-word": "docx",
        "word-to-pdf": "pdf", "excel-to-pdf": "pdf", "pptx-to-pdf": "pdf",
        "png-to-jpg": "jpg", "jpg-to-png": "png",
        "epub-to-mobi": "mobi", "mobi-to-epub": "epub", "epub-to-pdf": "pdf",
    }
    target_ext = ext_map.get(conversion_type, "bin")

    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for fname, file_data in file_buffers:
            base_name = os.path.splitext(fname)[0]
            try:
                response = _convert_single(file_data, base_name, conversion_type)
                # Read the streaming response body
                content = b""
                async for chunk in response.body_iterator:
                    if isinstance(chunk, str):
                        content += chunk.encode()
                    else:
                        content += chunk

                # PDF→images may return a zip — handle nested zip
                disposition = response.headers.get("content-disposition", "")
                if ".zip" in disposition:
                    inner_zip = zipfile.ZipFile(io.BytesIO(content))
                    for inner_name in inner_zip.namelist():
                        zf.writestr(inner_name, inner_zip.read(inner_name))
                    inner_zip.close()
                else:
                    zf.writestr(f"{base_name}.{target_ext}", content)
            except Exception as e:
                logger.warning(f"Batch: skipping {fname}: {e}")
                zf.writestr(f"{base_name}_ERROR.txt", f"Conversion failed: {str(e)}")

    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={"Content-Disposition": 'attachment; filename="batch_converted.zip"'},
    )


def _convert_single(file_data: bytes, name: str, conversion_type: str):
    """Route a single file to the correct converter, return a StreamingResponse."""
    if conversion_type == "pdf-to-jpg":
        return _pdf_to_images(file_data, name, "jpeg")
    elif conversion_type == "pdf-to-png":
        return _pdf_to_images(file_data, name, "png")
    elif conversion_type == "pdf-to-word":
        return _pdf_to_word(file_data, name)
    elif conversion_type in ("jpg-to-pdf", "png-to-pdf", "image-to-pdf"):
        return _image_to_pdf(file_data, name)
    elif conversion_type == "word-to-pdf":
        return _word_to_pdf(file_data, name)
    elif conversion_type == "excel-to-pdf":
        return _excel_to_pdf(file_data, name)
    elif conversion_type == "pptx-to-pdf":
        return _pptx_to_pdf(file_data, name)
    elif conversion_type == "png-to-jpg":
        return _convert_image_format(file_data, name, "jpeg")
    elif conversion_type == "jpg-to-png":
        return _convert_image_format(file_data, name, "png")
    elif conversion_type == "epub-to-mobi":
        return _calibre_convert(file_data, name, ".epub", ".mobi", "application/x-mobipocket-ebook")
    elif conversion_type == "mobi-to-epub":
        return _calibre_convert(file_data, name, ".mobi", ".epub", "application/epub+zip")
    elif conversion_type == "epub-to-pdf":
        return _epub_to_pdf(file_data, name)
    raise HTTPException(status_code=400, detail="Conversion not implemented")


# ── PDF to Images ──────────────────────────────────────────────
def _pdf_to_images(pdf_data: bytes, name: str, fmt: str):
    doc = fitz.open(stream=pdf_data, filetype="pdf")
    ext = "jpg" if fmt == "jpeg" else "png"
    if len(doc) == 1:
        pix = doc[0].get_pixmap(dpi=200)
        img_data = pix.tobytes(fmt)
        doc.close()
        return StreamingResponse(
            io.BytesIO(img_data),
            media_type=f"image/{fmt}",
            headers={"Content-Disposition": f'attachment; filename="{name}.{ext}"'},
        )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=200)
            zf.writestr(f"{name}_page_{i+1}.{ext}", pix.tobytes(fmt))
    doc.close()
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{name}_images.zip"'},
    )


# ── PDF to Word ────────────────────────────────────────────────
def _pdf_to_word(pdf_data: bytes, name: str):
    from pdf2docx import Converter
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
        tmp_pdf.write(pdf_data)
        tmp_pdf_path = tmp_pdf.name
    tmp_docx_path = tmp_pdf_path.replace(".pdf", ".docx")
    try:
        cv = Converter(tmp_pdf_path)
        cv.convert(tmp_docx_path)
        cv.close()
        with open(tmp_docx_path, "rb") as f:
            docx_data = f.read()
    finally:
        for p in (tmp_pdf_path, tmp_docx_path):
            if os.path.exists(p):
                os.unlink(p)
    return StreamingResponse(
        io.BytesIO(docx_data),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{name}.docx"'},
    )


# ── Image to PDF ──────────────────────────────────────────────
def _image_to_pdf(img_data: bytes, name: str):
    doc = fitz.open()
    img = fitz.open(stream=img_data, filetype="png")
    pdfbytes = img.convert_to_pdf()
    img_pdf = fitz.open("pdf", pdfbytes)
    doc.insert_pdf(img_pdf)
    pdf_bytes = doc.tobytes()
    doc.close()
    img.close()
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{name}.pdf"'},
    )


# ── Word to PDF ───────────────────────────────────────────────
def _word_to_pdf(docx_data: bytes, name: str):
    from docx import Document
    doc = Document(io.BytesIO(docx_data))
    pdf = fitz.open()
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        style = para.style.name if para.style else "Normal"
        lines.append((text, style))
    if not lines:
        lines = [("(Empty document)", "Normal")]

    max_lines = 45
    i = 0
    while i < len(lines):
        page = pdf.new_page(width=612, height=792)
        y, count = 50, 0
        while i < len(lines) and count < max_lines:
            text, style = lines[i]
            fs = 18 if "Heading 1" in style else 15 if "Heading 2" in style else 13 if "Heading" in style else 11
            if text:
                for wrapped in _wrap_text(text, fs, 500):
                    page.insert_text(fitz.Point(50, y), wrapped, fontsize=fs, color=(0.1, 0.1, 0.1))
                    y += fs + 4
                    count += 1
                    if count >= max_lines:
                        break
                y += 2
            else:
                y += 8
                count += 1
            i += 1
    pdf_bytes = pdf.tobytes()
    pdf.close()
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{name}.pdf"'},
    )


def _wrap_text(text, fontsize, max_width):
    words = text.split()
    lines, line = [], ""
    for word in words:
        test = f"{line} {word}".strip()
        if len(test) * fontsize * 0.45 > max_width:
            if line:
                lines.append(line)
            line = word
        else:
            line = test
    if line:
        lines.append(line)
    return lines or [""]


# ── Excel to PDF (Improved — proper table grid) ───────────────
def _excel_to_pdf(xlsx_data: bytes, name: str):
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(xlsx_data), data_only=True)
    pdf = fitz.open()
    PAGE_W, PAGE_H = 842, 595  # landscape A4
    MARGIN = 35
    HEADER_H = 30
    ROW_H = 16
    FONTSIZE = 8
    HEADER_FS = 9

    for ws in wb.worksheets:
        # Determine column count and widths
        max_col = min(ws.max_column or 1, 12)
        max_row = min(ws.max_row or 1, 200)
        usable_w = PAGE_W - 2 * MARGIN
        col_w = usable_w / max_col

        page = pdf.new_page(width=PAGE_W, height=PAGE_H)
        # Sheet title
        page.insert_text(fitz.Point(MARGIN, 25), ws.title, fontsize=12, color=(0.2, 0.15, 0.5))
        y = 38
        rows_on_page = 0

        for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col, values_only=False)):
            if y + ROW_H > PAGE_H - MARGIN:
                page = pdf.new_page(width=PAGE_W, height=PAGE_H)
                y = MARGIN
                rows_on_page = 0

            is_header = (row_idx == 0)
            cell_h = HEADER_H if is_header else ROW_H
            fs = HEADER_FS if is_header else FONTSIZE

            # Draw row background
            if is_header:
                rect = fitz.Rect(MARGIN, y, MARGIN + usable_w, y + cell_h)
                page.draw_rect(rect, color=None, fill=(0.22, 0.16, 0.52))
            elif row_idx % 2 == 0:
                rect = fitz.Rect(MARGIN, y, MARGIN + usable_w, y + cell_h)
                page.draw_rect(rect, color=None, fill=(0.96, 0.96, 0.98))

            # Draw cells
            for col_idx, cell in enumerate(row):
                x = MARGIN + col_idx * col_w
                cell_rect = fitz.Rect(x, y, x + col_w, y + cell_h)
                # Cell border
                page.draw_rect(cell_rect, color=(0.8, 0.8, 0.85), width=0.3)
                # Cell text
                val = str(cell.value) if cell.value is not None else ""
                if len(val) > 18:
                    val = val[:16] + ".."
                text_color = (1, 1, 1) if is_header else (0.15, 0.15, 0.15)
                page.insert_text(
                    fitz.Point(x + 3, y + cell_h - 4 if is_header else y + cell_h - 3),
                    val, fontsize=fs, color=text_color,
                )

            y += cell_h
            rows_on_page += 1

    pdf_bytes = pdf.tobytes()
    pdf.close()
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{name}.pdf"'},
    )


# ── PPTX to PDF (with image support) ─────────────────────────
def _pptx_to_pdf(pptx_data: bytes, name: str):
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(pptx_data))
    slide_w = prs.slide_width or Emu(9144000)
    slide_h = prs.slide_height or Emu(6858000)
    w_pt = slide_w / 12700
    h_pt = slide_h / 12700

    pdf = fitz.open()

    for slide_num, slide in enumerate(prs.slides):
        page = pdf.new_page(width=w_pt, height=h_pt)

        # Light background
        page.draw_rect(fitz.Rect(0, 0, w_pt, h_pt), color=None, fill=(0.98, 0.98, 1.0))

        # Slide number footer
        page.insert_text(
            fitz.Point(w_pt - 70, h_pt - 15),
            f"Slide {slide_num + 1}/{len(prs.slides)}",
            fontsize=7, color=(0.6, 0.6, 0.6),
        )

        y_text = 40

        for shape in slide.shapes:
            # Handle images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (hasattr(shape, "image") and shape.image):
                try:
                    img_blob = shape.image.blob
                    # Convert shape position to points
                    sx = shape.left / 12700 if shape.left else 40
                    sy = shape.top / 12700 if shape.top else y_text
                    sw = shape.width / 12700 if shape.width else 200
                    sh = shape.height / 12700 if shape.height else 150

                    # Clamp to page bounds
                    sx = max(10, min(sx, w_pt - 50))
                    sy = max(10, min(sy, h_pt - 50))
                    sw = min(sw, w_pt - sx - 10)
                    sh = min(sh, h_pt - sy - 10)

                    img_rect = fitz.Rect(sx, sy, sx + sw, sy + sh)
                    page.insert_image(img_rect, stream=img_blob)
                    y_text = max(y_text, sy + sh + 10)
                except Exception as img_err:
                    logger.debug(f"Could not render PPTX image: {img_err}")

            # Handle text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    fs = 11
                    color = (0.1, 0.1, 0.1)
                    if para.level == 0 and len(text) < 80:
                        fs = 18
                        color = (0.2, 0.15, 0.5)
                    elif para.level == 1:
                        fs = 14

                    for wrapped in _wrap_text(text, fs, w_pt - 80):
                        if y_text > h_pt - 40:
                            break
                        page.insert_text(fitz.Point(40, y_text), wrapped, fontsize=fs, color=color)
                        y_text += fs + 6

                    if y_text > h_pt - 40:
                        break

    pdf_bytes = pdf.tobytes()
    pdf.close()
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{name}.pdf"'},
    )


# ── EPUB to PDF (native — ebooklib + PyMuPDF) ────────────────
def _epub_to_pdf(epub_data: bytes, name: str):
    from ebooklib import epub
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.chunks = []
            self._tag_stack = []

        def handle_starttag(self, tag, attrs):
            self._tag_stack.append(tag)

        def handle_endtag(self, tag):
            if self._tag_stack and self._tag_stack[-1] == tag:
                self._tag_stack.pop()
            if tag in ("p", "div", "br", "li"):
                self.chunks.append(("", "break"))

        def handle_data(self, data):
            text = data.strip()
            if not text:
                return
            tag = self._tag_stack[-1] if self._tag_stack else "p"
            if tag == "h1":
                self.chunks.append((text, "h1"))
            elif tag == "h2":
                self.chunks.append((text, "h2"))
            elif tag in ("h3", "h4", "h5", "h6"):
                self.chunks.append((text, "h3"))
            else:
                self.chunks.append((text, "body"))

    book = epub.read_epub(io.BytesIO(epub_data))
    all_chunks = []

    # Book title
    title = book.get_metadata("DC", "title")
    if title:
        all_chunks.append((title[0][0], "title"))
        all_chunks.append(("", "break"))

    for item in book.get_items_of_type(9):  # ITEM_DOCUMENT
        content = item.get_content().decode("utf-8", errors="ignore")
        parser = _TextExtractor()
        parser.feed(content)
        all_chunks.extend(parser.chunks)
        all_chunks.append(("", "break"))

    if not all_chunks:
        all_chunks = [("(Empty eBook)", "body")]

    # Render to PDF
    pdf = fitz.open()
    PAGE_W, PAGE_H = 432, 648  # 6x9 inches — book format
    MARGIN_X, MARGIN_TOP, MARGIN_BOT = 50, 50, 50
    max_w = PAGE_W - 2 * MARGIN_X

    page = pdf.new_page(width=PAGE_W, height=PAGE_H)
    y = MARGIN_TOP

    for text, style in all_chunks:
        if style == "break":
            y += 8
            if y > PAGE_H - MARGIN_BOT:
                page = pdf.new_page(width=PAGE_W, height=PAGE_H)
                y = MARGIN_TOP
            continue

        if style == "title":
            fs, color = 22, (0.2, 0.12, 0.45)
        elif style == "h1":
            fs, color = 17, (0.2, 0.15, 0.5)
            y += 8
        elif style == "h2":
            fs, color = 14, (0.25, 0.2, 0.5)
            y += 4
        elif style == "h3":
            fs, color = 12, (0.3, 0.25, 0.45)
            y += 2
        else:
            fs, color = 10, (0.12, 0.12, 0.12)

        for line in _wrap_text(text, fs, max_w):
            if y + fs + 4 > PAGE_H - MARGIN_BOT:
                page = pdf.new_page(width=PAGE_W, height=PAGE_H)
                y = MARGIN_TOP
            page.insert_text(fitz.Point(MARGIN_X, y), line, fontsize=fs, color=color)
            y += fs + 3

    pdf_bytes = pdf.tobytes()
    pdf.close()
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{name}.pdf"'},
    )


# ── Calibre-based conversions (EPUB/MOBI) ─────────────────────
def _calibre_convert(file_data: bytes, name: str, src_ext: str, dst_ext: str, mime: str):
    with tempfile.TemporaryDirectory() as tmpdir:
        src_path = os.path.join(tmpdir, f"input{src_ext}")
        dst_path = os.path.join(tmpdir, f"output{dst_ext}")

        with open(src_path, "wb") as f:
            f.write(file_data)

        # PDF output requires a virtual display for Qt WebEngine
        cmd = ["ebook-convert", src_path, dst_path]
        if dst_ext == ".pdf":
            cmd = ["xvfb-run", "--auto-servernum", "--server-args=-screen 0 1024x768x24"] + cmd

        result = subprocess.run(
            cmd,
            capture_output=True, text=True, timeout=120,
        )

        if result.returncode != 0:
            logger.error(f"Calibre error: {result.stderr[:500]}")
            raise HTTPException(status_code=500, detail=f"eBook conversion failed: {result.stderr[:200]}")

        if not os.path.exists(dst_path):
            raise HTTPException(status_code=500, detail="Conversion produced no output")

        with open(dst_path, "rb") as f:
            out_data = f.read()

    ext = dst_ext.lstrip(".")
    return StreamingResponse(
        io.BytesIO(out_data),
        media_type=mime,
        headers={"Content-Disposition": f'attachment; filename="{name}.{ext}"'},
    )


# ── Image format conversion ───────────────────────────────────
def _convert_image_format(img_data: bytes, name: str, target_fmt: str):
    img = Image.open(io.BytesIO(img_data))
    if img.mode == "RGBA" and target_fmt == "jpeg":
        img = img.convert("RGB")
    buf = io.BytesIO()
    ext = "jpg" if target_fmt == "jpeg" else "png"
    img.save(buf, format=target_fmt.upper(), quality=92)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type=f"image/{target_fmt}",
        headers={"Content-Disposition": f'attachment; filename="{name}.{ext}"'},
    )
