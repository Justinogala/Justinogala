"""
Presentations routes - CRUD for slide presentations.
"""
from fastapi import APIRouter, HTTPException, Depends, Query
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from datetime import datetime, timezone
from typing import Optional
import uuid
import jwt
import os

from config import db, JWT_SECRET_KEY, JWT_ALGORITHM, logger

router = APIRouter(prefix="/presentations", tags=["Presentations"])
security = HTTPBearer(auto_error=False)


def _get_user_id(creds: HTTPAuthorizationCredentials):
    if not creds:
        raise HTTPException(status_code=401, detail="Not authenticated")
    try:
        payload = jwt.decode(creds.credentials, JWT_SECRET_KEY, algorithms=[JWT_ALGORITHM])
        return payload.get("user_id") or payload.get("sub")
    except jwt.PyJWTError:
        raise HTTPException(status_code=401, detail="Invalid token")


DEFAULT_SLIDES = [
    {"id": "slide-1", "layout": "title", "title": "Untitled Presentation", "subtitle": "Click to add subtitle", "notes": ""},
]


@router.get("")
async def list_presentations(
    search: Optional[str] = Query(None),
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    user_id = _get_user_id(credentials)
    query = {"user_id": user_id, "deleted": {"$ne": True}}
    if search:
        query["title"] = {"$regex": search, "$options": "i"}
    docs = await db.presentations.find(query, {"_id": 0, "slides": 0}).sort("updated_at", -1).to_list(200)
    return docs


@router.post("")
async def create_presentation(
    body: dict,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    user_id = _get_user_id(credentials)
    pres_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    slides = body.get("slides", DEFAULT_SLIDES)

    doc = {
        "id": pres_id,
        "user_id": user_id,
        "title": body.get("title", "Untitled Presentation"),
        "slides": slides,
        "slide_count": len(slides),
        "template": body.get("template"),
        "created_at": now,
        "updated_at": now,
        "deleted": False,
    }
    await db.presentations.insert_one(doc)
    doc.pop("_id", None)
    return doc


@router.get("/{pres_id}")
async def get_presentation(
    pres_id: str,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    user_id = _get_user_id(credentials)
    doc = await db.presentations.find_one({"id": pres_id, "user_id": user_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Presentation not found")
    return doc


@router.put("/{pres_id}")
async def update_presentation(
    pres_id: str,
    body: dict,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    user_id = _get_user_id(credentials)
    existing = await db.presentations.find_one({"id": pres_id, "user_id": user_id})
    if not existing:
        raise HTTPException(status_code=404, detail="Presentation not found")

    update_data = {"updated_at": datetime.now(timezone.utc).isoformat()}
    if "title" in body:
        update_data["title"] = body["title"]
    if "slides" in body:
        update_data["slides"] = body["slides"]
        update_data["slide_count"] = len(body["slides"])

    await db.presentations.update_one({"id": pres_id}, {"$set": update_data})
    updated = await db.presentations.find_one({"id": pres_id}, {"_id": 0})
    return updated


@router.delete("/{pres_id}")
async def delete_presentation(
    pres_id: str,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    user_id = _get_user_id(credentials)
    result = await db.presentations.update_one(
        {"id": pres_id, "user_id": user_id},
        {"$set": {"deleted": True, "updated_at": datetime.now(timezone.utc).isoformat()}}
    )
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="Presentation not found")
    return {"success": True}


@router.post("/{pres_id}/duplicate")
async def duplicate_presentation(
    pres_id: str,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    user_id = _get_user_id(credentials)
    original = await db.presentations.find_one({"id": pres_id, "user_id": user_id}, {"_id": 0})
    if not original:
        raise HTTPException(status_code=404, detail="Presentation not found")

    new_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    new_doc = {**original, "id": new_id, "title": f"{original['title']} (Copy)", "created_at": now, "updated_at": now}
    await db.presentations.insert_one(new_doc)
    new_doc.pop("_id", None)
    return new_doc


@router.post("/ai-generate")
async def ai_generate_presentation(
    body: dict,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    """Generate a presentation using AI"""
    from emergentintegrations.llm.chat import LlmChat, UserMessage
    import json as json_mod

    user_id = _get_user_id(credentials)
    prompt = body.get("prompt", "").strip()
    if not prompt:
        raise HTTPException(status_code=400, detail="Prompt is required")

    api_key = os.environ.get("EMERGENT_LLM_KEY")
    if not api_key:
        raise HTTPException(status_code=503, detail="AI service not configured")

    system_msg = """You are a presentation designer. Generate a JSON array of slides for a professional presentation.
Each slide object must have: "layout" (one of: "title", "content", "two-column", "section", "blank"), "title" (string), and optional fields based on layout:
- "title" layout: "subtitle" (string)
- "content" layout: "body" (string, use bullet points separated by newlines)
- "two-column" layout: "left" (string), "right" (string)
- "section" layout: "subtitle" (string)
Generate 6-10 slides. Return ONLY valid JSON array, no markdown fences."""

    try:
        chat = LlmChat(
            api_key=api_key,
            session_id=f"pres-gen-{uuid.uuid4()}",
            system_message=system_msg
        ).with_model("openai", "gpt-5.2")

        user_message = UserMessage(text=f"Create a presentation about: {prompt}")
        response = await chat.send_message(user_message)

        # Parse JSON from response
        response = response.strip()
        if response.startswith("```"):
            response = response.split("\n", 1)[1].rsplit("```", 1)[0]
        slides = json_mod.loads(response)

        # Add IDs to slides
        for i, slide in enumerate(slides):
            slide["id"] = f"slide-{i+1}"
            slide["notes"] = ""

        title = slides[0].get("title", prompt[:60]) if slides else prompt[:60]

    except Exception as e:
        logger.error(f"AI presentation generation failed: {e}")
        # Fallback slides
        title = prompt[:60]
        slides = [
            {"id": "slide-1", "layout": "title", "title": title, "subtitle": "AI-Generated Presentation", "notes": ""},
            {"id": "slide-2", "layout": "content", "title": "Overview", "body": "Key points about this topic\nAdd your content here", "notes": ""},
            {"id": "slide-3", "layout": "content", "title": "Details", "body": "Supporting information\nData and evidence", "notes": ""},
            {"id": "slide-4", "layout": "content", "title": "Next Steps", "body": "Action item 1\nAction item 2\nAction item 3", "notes": ""},
            {"id": "slide-5", "layout": "section", "title": "Thank You", "subtitle": "Questions?", "notes": ""},
        ]

    pres_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    doc = {
        "id": pres_id,
        "user_id": user_id,
        "title": title,
        "slides": slides,
        "slide_count": len(slides),
        "template": "ai-generated",
        "created_at": now,
        "updated_at": now,
        "deleted": False,
    }
    await db.presentations.insert_one(doc)
    doc.pop("_id", None)
    return doc


@router.get("/{pres_id}/export/pptx")
async def export_pptx(
    pres_id: str,
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    """Export presentation as PPTX"""
    from fastapi.responses import Response
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from io import BytesIO
    import re

    user_id = _get_user_id(credentials)
    doc = await db.presentations.find_one({"id": pres_id, "user_id": user_id}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Presentation not found")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in doc.get("slides", []):
        layout_name = slide_data.get("layout", "content")

        if layout_name == "title":
            slide_layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            if slide.placeholders[1]:
                slide.placeholders[1].text = slide_data.get("subtitle", "")

        elif layout_name == "section":
            slide_layout = prs.slide_layouts[2]  # Section Header
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get("subtitle", "")

        elif layout_name == "two-column":
            slide_layout = prs.slide_layouts[3]  # Two Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data.get("left", "")
            if len(slide.placeholders) > 2:
                slide.placeholders[2].text = slide_data.get("right", "")

        elif layout_name == "blank":
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

        else:  # content
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            body = slide_data.get("body", "")
            if body and len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, line in enumerate(body.split("\n")):
                    if i == 0:
                        tf.paragraphs[0].text = line.strip()
                    else:
                        p = tf.add_paragraph()
                        p.text = line.strip()

        # Add notes
        notes = slide_data.get("notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    filename = re.sub(r'[^\w\s\-]', '', doc.get("title", "presentation")).strip() or "presentation"

    return Response(
        content=buffer.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}.pptx"'}
    )
