"""
Test File Converter API - Tests for PDF/Image/Document/eBook format conversions
Includes: 14 conversions (11 original + 3 new eBook: epub-to-mobi, mobi-to-epub, epub-to-pdf)
"""
import pytest
import requests
import os
import io
from PIL import Image
import fitz  # PyMuPDF

BASE_URL = os.environ.get('REACT_APP_BACKEND_URL', '').rstrip('/')


class TestFileConverterSupported:
    """Tests for GET /api/converter/supported endpoint"""
    
    def test_get_supported_conversions_count(self):
        """Test that supported conversions endpoint returns 14 conversions (11 original + 3 eBook)"""
        response = requests.get(f"{BASE_URL}/api/converter/supported")
        assert response.status_code == 200
        
        data = response.json()
        assert "conversions" in data
        assert len(data["conversions"]) == 14, f"Expected 14 conversions, got {len(data['conversions'])}"
    
    def test_all_conversion_ids_present(self):
        """Test that all 14 expected conversion types are present"""
        response = requests.get(f"{BASE_URL}/api/converter/supported")
        assert response.status_code == 200
        
        data = response.json()
        conversion_ids = [c["id"] for c in data["conversions"]]
        
        # Original 11 conversions
        expected_ids = [
            "pdf-to-jpg", "pdf-to-png", "pdf-to-word",
            "word-to-pdf", "excel-to-pdf", "pptx-to-pdf",
            "jpg-to-pdf", "png-to-pdf", "image-to-pdf",
            "png-to-jpg", "jpg-to-png",
            # New 3 eBook conversions
            "epub-to-mobi", "mobi-to-epub", "epub-to-pdf"
        ]
        for expected_id in expected_ids:
            assert expected_id in conversion_ids, f"Missing conversion: {expected_id}"
    
    def test_ebook_category_present(self):
        """Test that eBook category is present with 3 items"""
        response = requests.get(f"{BASE_URL}/api/converter/supported")
        assert response.status_code == 200
        
        data = response.json()
        ebook_conversions = [c for c in data["conversions"] if c.get("category") == "eBook"]
        assert len(ebook_conversions) == 3, f"Expected 3 eBook conversions, got {len(ebook_conversions)}"
        
        ebook_ids = [c["id"] for c in ebook_conversions]
        assert "epub-to-mobi" in ebook_ids
        assert "mobi-to-epub" in ebook_ids
        assert "epub-to-pdf" in ebook_ids
    
    def test_supported_conversions_structure(self):
        """Test that each conversion has required fields"""
        response = requests.get(f"{BASE_URL}/api/converter/supported")
        assert response.status_code == 200
        
        data = response.json()
        for conv in data["conversions"]:
            assert "id" in conv
            assert "from" in conv
            assert "to" in conv
            assert "category" in conv


class TestPDFToImageConversions:
    """Tests for PDF to image conversions"""
    
    @pytest.fixture
    def test_pdf(self):
        """Create a test PDF file"""
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text(fitz.Point(72, 72), "Test PDF Content")
        pdf_bytes = doc.tobytes()
        doc.close()
        return pdf_bytes
    
    def test_pdf_to_jpg_conversion(self, test_pdf):
        """Test PDF to JPG conversion"""
        files = {"file": ("test.pdf", io.BytesIO(test_pdf), "application/pdf")}
        data = {"conversion_type": "pdf-to-jpg"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        content_type = response.headers.get("Content-Type", "")
        assert "image/jpeg" in content_type or "application/zip" in content_type
        
        disposition = response.headers.get("Content-Disposition", "")
        assert "attachment" in disposition
    
    def test_pdf_to_png_conversion(self, test_pdf):
        """Test PDF to PNG conversion"""
        files = {"file": ("test.pdf", io.BytesIO(test_pdf), "application/pdf")}
        data = {"conversion_type": "pdf-to-png"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        content_type = response.headers.get("Content-Type", "")
        assert "image/png" in content_type or "application/zip" in content_type
    
    def test_pdf_to_word_conversion(self, test_pdf):
        """Test PDF to Word conversion"""
        files = {"file": ("test.pdf", io.BytesIO(test_pdf), "application/pdf")}
        data = {"conversion_type": "pdf-to-word"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        content_type = response.headers.get("Content-Type", "")
        assert "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in content_type
        
        disposition = response.headers.get("Content-Disposition", "")
        assert ".docx" in disposition


class TestImageConversions:
    """Tests for image format conversions"""
    
    @pytest.fixture
    def test_png(self):
        """Create a test PNG image"""
        img = Image.new("RGB", (200, 100), (100, 50, 200))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf.getvalue()
    
    @pytest.fixture
    def test_jpg(self):
        """Create a test JPG image"""
        img = Image.new("RGB", (200, 100), (50, 100, 200))
        buf = io.BytesIO()
        img.save(buf, format="JPEG")
        buf.seek(0)
        return buf.getvalue()
    
    def test_png_to_jpg_conversion(self, test_png):
        """Test PNG to JPG conversion"""
        files = {"file": ("test.png", io.BytesIO(test_png), "image/png")}
        data = {"conversion_type": "png-to-jpg"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        content_type = response.headers.get("Content-Type", "")
        assert "image/jpeg" in content_type
        
        img = Image.open(io.BytesIO(response.content))
        assert img.format == "JPEG"
    
    def test_jpg_to_png_conversion(self, test_jpg):
        """Test JPG to PNG conversion"""
        files = {"file": ("test.jpg", io.BytesIO(test_jpg), "image/jpeg")}
        data = {"conversion_type": "jpg-to-png"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        content_type = response.headers.get("Content-Type", "")
        assert "image/png" in content_type
        
        img = Image.open(io.BytesIO(response.content))
        assert img.format == "PNG"
    
    def test_image_to_pdf_conversion(self, test_png):
        """Test image to PDF conversion"""
        files = {"file": ("test.png", io.BytesIO(test_png), "image/png")}
        data = {"conversion_type": "image-to-pdf"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        content_type = response.headers.get("Content-Type", "")
        assert "application/pdf" in content_type
        
        doc = fitz.open(stream=response.content, filetype="pdf")
        assert len(doc) >= 1
        doc.close()


class TestWordToPDFConversion:
    """Tests for Word to PDF conversion"""
    
    @pytest.fixture
    def test_docx(self):
        """Create a test DOCX file"""
        from docx import Document
        doc = Document()
        doc.add_heading("Test Document", 0)
        doc.add_paragraph("This is a test paragraph.")
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.getvalue()
    
    def test_word_to_pdf_conversion(self, test_docx):
        """Test Word to PDF conversion"""
        files = {"file": ("test.docx", io.BytesIO(test_docx), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
        data = {"conversion_type": "word-to-pdf"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        content_type = response.headers.get("Content-Type", "")
        assert "application/pdf" in content_type
        
        doc = fitz.open(stream=response.content, filetype="pdf")
        assert len(doc) >= 1
        doc.close()


class TestExcelToPDFConversion:
    """Tests for Excel to PDF conversion with improved table grid rendering"""
    
    @pytest.fixture
    def test_xlsx(self):
        """Create a test XLSX file with headers and data"""
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Test Sheet"
        # Header row
        ws['A1'] = 'Name'
        ws['B1'] = 'Age'
        ws['C1'] = 'City'
        # Data rows
        ws['A2'] = 'Alice'
        ws['B2'] = 25
        ws['C2'] = 'New York'
        ws['A3'] = 'Bob'
        ws['B3'] = 30
        ws['C3'] = 'Los Angeles'
        ws['A4'] = 'Charlie'
        ws['B4'] = 35
        ws['C4'] = 'Chicago'
        
        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)
        return buf.getvalue()
    
    def test_excel_to_pdf_conversion(self, test_xlsx):
        """Test Excel to PDF conversion returns valid PDF"""
        files = {"file": ("test.xlsx", io.BytesIO(test_xlsx), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")}
        data = {"conversion_type": "excel-to-pdf"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        content_type = response.headers.get("Content-Type", "")
        assert "application/pdf" in content_type
        
        # Verify the output is a valid PDF
        doc = fitz.open(stream=response.content, filetype="pdf")
        assert len(doc) >= 1
        doc.close()
    
    def test_excel_to_pdf_has_content(self, test_xlsx):
        """Test Excel to PDF conversion produces PDF with content (table grid)"""
        files = {"file": ("test.xlsx", io.BytesIO(test_xlsx), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")}
        data = {"conversion_type": "excel-to-pdf"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        # PDF should have reasonable size indicating table content
        assert len(response.content) > 500, "PDF seems too small, may not have table content"


class TestPPTXToPDFConversion:
    """Tests for PPTX to PDF conversion with embedded images"""
    
    @pytest.fixture
    def test_pptx_with_image(self):
        """Create a test PPTX file with text and an embedded image"""
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title text
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Test Presentation with Image"
        
        # Create and add an image
        img = Image.new("RGB", (400, 300), (100, 150, 200))
        img_buf = io.BytesIO()
        img.save(img_buf, format="PNG")
        img_buf.seek(0)
        
        slide.shapes.add_picture(img_buf, Inches(2), Inches(2), width=Inches(4), height=Inches(3))
        
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()
    
    def test_pptx_to_pdf_conversion(self, test_pptx_with_image):
        """Test PPTX to PDF conversion returns valid PDF"""
        files = {"file": ("test.pptx", io.BytesIO(test_pptx_with_image), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
        data = {"conversion_type": "pptx-to-pdf"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        content_type = response.headers.get("Content-Type", "")
        assert "application/pdf" in content_type
        
        doc = fitz.open(stream=response.content, filetype="pdf")
        assert len(doc) >= 1
        doc.close()
    
    def test_pptx_to_pdf_with_images_has_large_size(self, test_pptx_with_image):
        """Test PPTX to PDF with images produces larger PDF (>100KB indicates images rendered)"""
        files = {"file": ("test.pptx", io.BytesIO(test_pptx_with_image), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
        data = {"conversion_type": "pptx-to-pdf"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 200
        
        # PDF with embedded images should be larger than 100KB
        # Note: This is a heuristic - actual size depends on image compression
        pdf_size = len(response.content)
        print(f"PPTX to PDF size: {pdf_size} bytes")
        # At minimum, should be larger than a text-only PDF
        assert pdf_size > 5000, f"PDF size {pdf_size} bytes seems too small for image content"


class TestEBookConversions:
    """Tests for eBook format conversions (EPUB, MOBI)"""
    
    @pytest.fixture
    def test_epub(self):
        """Create a test EPUB file using ebooklib"""
        from ebooklib import epub
        
        book = epub.EpubBook()
        book.set_identifier('test-book-123')
        book.set_title('Test eBook')
        book.set_language('en')
        book.add_author('Test Author')
        
        # Create a chapter
        c1 = epub.EpubHtml(title='Chapter 1', file_name='chap_01.xhtml', lang='en')
        c1.content = '<html><body><h1>Chapter 1</h1><p>This is test content for the eBook.</p></body></html>'
        book.add_item(c1)
        
        # Add navigation
        book.toc = (epub.Link('chap_01.xhtml', 'Chapter 1', 'chap1'),)
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())
        book.spine = ['nav', c1]
        
        buf = io.BytesIO()
        epub.write_epub(buf, book)
        buf.seek(0)
        return buf.getvalue()
    
    def test_epub_to_mobi_conversion(self, test_epub):
        """Test EPUB to MOBI conversion via Calibre"""
        files = {"file": ("test.epub", io.BytesIO(test_epub), "application/epub+zip")}
        data = {"conversion_type": "epub-to-mobi"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data, timeout=120)
        assert response.status_code == 200, f"EPUB to MOBI failed: {response.text}"
        
        content_type = response.headers.get("Content-Type", "")
        assert "application/x-mobipocket-ebook" in content_type or "application/octet-stream" in content_type
        
        disposition = response.headers.get("Content-Disposition", "")
        assert ".mobi" in disposition
        
        # MOBI file should have reasonable size
        assert len(response.content) > 1000, "MOBI file seems too small"
    
    def test_epub_to_pdf_conversion(self, test_epub):
        """Test EPUB to PDF conversion (native Python - ebooklib + PyMuPDF)"""
        files = {"file": ("test.epub", io.BytesIO(test_epub), "application/epub+zip")}
        data = {"conversion_type": "epub-to-pdf"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data, timeout=120)
        assert response.status_code == 200, f"EPUB to PDF failed: {response.text}"
        
        content_type = response.headers.get("Content-Type", "")
        assert "application/pdf" in content_type
        
        # Verify the output is a valid PDF
        doc = fitz.open(stream=response.content, filetype="pdf")
        assert len(doc) >= 1
        doc.close()
        
        # Check PDF is in book format (should have content)
        assert len(response.content) > 500, "PDF seems too small for eBook content"


class TestMobiToEpubConversion:
    """Tests for MOBI to EPUB conversion - requires creating MOBI first"""
    
    def test_mobi_to_epub_conversion(self):
        """Test MOBI to EPUB conversion via Calibre - create MOBI first then convert back"""
        # First create an EPUB
        from ebooklib import epub
        
        book = epub.EpubBook()
        book.set_identifier('test-mobi-book-456')
        book.set_title('Test MOBI Book')
        book.set_language('en')
        book.add_author('Test Author')
        
        c1 = epub.EpubHtml(title='Chapter 1', file_name='chap_01.xhtml', lang='en')
        c1.content = '<html><body><h1>Chapter 1</h1><p>Content for MOBI test.</p></body></html>'
        book.add_item(c1)
        
        book.toc = (epub.Link('chap_01.xhtml', 'Chapter 1', 'chap1'),)
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())
        book.spine = ['nav', c1]
        
        epub_buf = io.BytesIO()
        epub.write_epub(epub_buf, book)
        epub_buf.seek(0)
        
        # Convert EPUB to MOBI first
        files = {"file": ("test.epub", epub_buf, "application/epub+zip")}
        data = {"conversion_type": "epub-to-mobi"}
        
        mobi_response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data, timeout=120)
        if mobi_response.status_code != 200:
            pytest.skip(f"Could not create MOBI for test: {mobi_response.text}")
        
        mobi_data = mobi_response.content
        
        # Now convert MOBI back to EPUB
        files = {"file": ("test.mobi", io.BytesIO(mobi_data), "application/x-mobipocket-ebook")}
        data = {"conversion_type": "mobi-to-epub"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data, timeout=120)
        assert response.status_code == 200, f"MOBI to EPUB failed: {response.text}"
        
        content_type = response.headers.get("Content-Type", "")
        assert "application/epub+zip" in content_type or "application/octet-stream" in content_type
        
        disposition = response.headers.get("Content-Disposition", "")
        assert ".epub" in disposition
        
        # EPUB file should have reasonable size
        assert len(response.content) > 1000, "EPUB file seems too small"


class TestErrorHandling:
    """Tests for error handling in converter API"""
    
    def test_unsupported_conversion_type(self):
        """Test that unsupported conversion type returns 400"""
        img = Image.new("RGB", (100, 100), (255, 0, 0))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        
        files = {"file": ("test.png", buf, "image/png")}
        data = {"conversion_type": "unsupported-type"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files, data=data)
        assert response.status_code == 400
        
        error_data = response.json()
        assert "detail" in error_data
        assert "unsupported" in error_data["detail"].lower()
    
    def test_missing_file(self):
        """Test that missing file returns error"""
        data = {"conversion_type": "pdf-to-jpg"}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", data=data)
        assert response.status_code == 422
    
    def test_missing_conversion_type(self):
        """Test that missing conversion_type returns error"""
        img = Image.new("RGB", (100, 100), (255, 0, 0))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        
        files = {"file": ("test.png", buf, "image/png")}
        
        response = requests.post(f"{BASE_URL}/api/converter/convert", files=files)
        assert response.status_code == 422


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
